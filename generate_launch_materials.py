import os
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PtxPt
from pptx.dml.color import RGBColor as PtxRGB

def create_docx(filename):
    doc = Document()
    
    # Title
    title = doc.add_heading('Big Launch: Beacon Programming Language', 0)
    
    # Intro
    p = doc.add_paragraph()
    runner = p.add_run("üöÄ I built my own programming language!")
    runner.bold = True
    runner.font.size = Pt(14)
    
    doc.add_paragraph("Hey everyone! I‚Äôm super excited to share Beacon (BPL) ‚Äì a project I‚Äôve been building from scratch. It‚Äôs a new coding language designed to be as readable as English but powered by a fast C engine. ‚ö°")
    
    # Features
    doc.add_heading('Why Beacon?', level=1)
    
    features = [
        "Readable Syntax: Designed to feel like natural language.",
        "High Performance: Powered by a custom C runtime.",
        "Smart Types: Automatic type detection (no more casting!).",
        "Modern Features: Pack collections, Traverse loops, and more."
    ]
    
    for f in features:
        doc.add_paragraph(f, style='List Bullet')
        
    doc.add_heading('New in v2.0.0', level=1)
    doc.add_paragraph("I just released version 2.0.0 today! It includes:")
    doc.add_paragraph("‚Ä¢ Pack Collections: pack(1, 2, 3)", style='List Bullet')
    doc.add_paragraph("‚Ä¢ Traverse Loops: traverse i from 1 to 10", style='List Bullet')
    
    # Links
    doc.add_heading('Check it out:', level=1)
    
    p = doc.add_paragraph()
    p.add_run("üåê Website: ").bold = True
    p.add_run("https://adividan1012.github.io/Beacon-Labs/")
    
    p = doc.add_paragraph()
    p.add_run("üí¨ Join WhatsApp Community: ").bold = True
    p.add_run("https://chat.whatsapp.com/DRWgaZFQUlELZNwqVqLsw1")
    
    doc.save(filename)
    print(f"Created {filename}")

def create_pptx(filename):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Beacon (BPL)"
    subtitle.text = "Project Launch v2.0.0\nCreated by Aaditya"
    
    # Slide 2: What is Beacon?
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "What is Beacon?"
    
    content = slide.placeholders[1]
    content.text = "A new programming language built from scratch."
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "Readable: Code looks like English sentences"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Fast: Powered by a custom C Runtime"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Smart: Automatic Type Detection"
    p.level = 1

    # Slide 3: New Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "New in v2.0"
    
    content = slide.placeholders[1]
    content.text = "Major updates released today:"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "Pack Collections"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "firm nums = pack(1, 2, 3)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Traverse Loops"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "traverse i from 1 to 10"
    p.level = 2

    # Slide 4: Join Us
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Try it Out!"
    
    content = slide.placeholders[1]
    content.text = "Download and join the discussion."
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "üåê Website: adividan1012.github.io/Beacon-Labs"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "üí¨ WhatsApp Group: Join via Link"
    p.level = 0
    
    prs.save(filename)
    print(f"Created {filename}")

if __name__ == "__main__":
    create_docx("Beacon_Launch_Announcement.docx")
    create_pptx("Beacon_Launch_Presentation.pptx")
